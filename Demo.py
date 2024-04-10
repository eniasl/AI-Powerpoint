from pptx import Presentation
import win32com.client
import os, pptx
import pythoncom

def get_slide_count(prs):
  """ Get the number of slides in PPTX presentation """
  slidecount = 0
  for slide in prs.slides:
    slidecount += 1
  return slidecount


def delete_slide(prs, slide):
  """ Delete a slide out of a powerpoint presentation"""
  id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
  slide_id = slide.slide_id
  prs.part.drop_rel(id_dict[slide_id][1])
  del prs.slides._sldIdLst[id_dict[slide_id][0]]


def get_single_slide_pres(prs, slidetokeep):
  for idx, slide in enumerate(prs.slides):
    if idx < slidetokeep:
      delete_slide(prs, slide)
    elif (idx > slidetokeep):
      delete_slide(prs, slide)
  prs.save("temp\\" + str(slidetokeep + 1) + ".pptx")

def mergePresentations(inputFileNames, outputFileName):
    os.system("TASKKILL /F /IM powerpnt.exe")
    Application = win32com.client.Dispatch("PowerPoint.Application",pythoncom.CoInitialize())
    outputPresentation = Application.Presentations.Add()
    outputPresentation.SaveAs(outputFileName)

    for file in inputFileNames:
        currentPresentation = Application.Presentations.Open(file)
        currentPresentation.Slides.Range(range(1, currentPresentation.Slides.Count+1)).copy()
        Application.Presentations(outputFileName).Windows(1).Activate()
        outputPresentation.Application.CommandBars.ExecuteMso("PasteSourceFormatting")
        currentPresentation.Close()

    outputPresentation.save()
    outputPresentation.close()
    Application.Quit()


def copy_slide(prs, slide_no):
    pptxfilepath = prs
    prs = Presentation(pptxfilepath)
    slidecount = get_slide_count(prs)
    for i in range(slidecount):
      prs_backup = Presentation(pptxfilepath)
      get_single_slide_pres(prs_backup, i)
      prs_backup = None

    # os.system("TASKKILL /F /IM powerpnt.exe")
    cwd = os.getcwd()
    lst = []
    dir_path = cwd+'\\temp\\'
    for path in os.listdir(dir_path):
        # check if current path is a file
        if os.path.isfile(os.path.join(dir_path, path)):
            lst.append(dir_path + path)

    lst.append(dir_path + str(slide_no) + ".pptx")
    print(lst)

    os.remove(os.getcwd() + "\\" + pptxfilepath)
    mergePresentations(lst, os.getcwd() + "\\" + pptxfilepath)

    for file_name in os.listdir(dir_path):
        file = dir_path + file_name
        if os.path.isfile(file):
            os.remove(file)


def update_slide_text(input_pptx, list_data):
    prs = pptx.Presentation(input_pptx)
    """get to the required slide"""
    slide = prs.slides[len(prs.slides)-1]
    """Find required text box"""
    for shape, data in zip(slide.shapes, list_data):
        if not shape.has_text_frame:
            continue
        shape.text_frame.text = data
    """save the file"""
    prs.save(input_pptx)


def create_pres(topic, mainpoints,slide_content, template):
    ppt = Presentation(template)
    ppt_name = "output/my_presentation.pptx"
    ppt.save(ppt_name)
    title_slide = [topic,"this is a test"]
    copy_slide(ppt_name, 2)
    update_slide_text(ppt_name,title_slide)
    copy_slide(ppt_name, 1)
    content_slide = ["content", mainpoints]
    update_slide_text(ppt_name, content_slide)
    for slide in slide_content:
        copy_slide(ppt_name, 1)
        update_slide_text(ppt_name, slide)
    copy_slide(ppt_name, 3)
    update_slide_text(ppt_name, ["thanks for listening","created by python"])

    os.system("TASKKILL /F /IM powerpnt.exe")


