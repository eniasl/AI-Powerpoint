from pptx import Presentation
import comtypes.client
import os

from pptx.dml.color import RGBColor
from pptx.util import Pt


def createPresentation(outputFileName, content, topic):
    root = Presentation()
    first_slide_layout = root.slide_layouts[0]
    second_slide_layout = root.slide_layouts[1]

    """ Ref for slide types:
    0 -> title and subtitle
    1 -> title and content
    2 -> section header
    3 -> two content
    4 -> Comparison
    5 -> Title only
    6 -> Blank
    7 -> Content with caption
    8 -> Pic with caption
    """

    # Creating slide object to add
    # in ppt i.e. Attaching slides
    # with Presentation i.e. ppt
    for row in content:
        slide = root.slides.add_slide(second_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,0,0)

        # Adding title and subtitle in
        # slide i.e. first page of slide
        slide.shapes.title.text = row[0]

        # We have different formats of
        # subtitles in ppts, for simple
        # subtitle this method should
        # implemented, you can change
        # 0 to 1 for different design
        text = ""
        check = True
        for col in row:
            if check:
                check = False
            else:
                text = text + col + "\n"

        slide.placeholders[1].text = text
        for paragraph in slide.placeholders[1].text_frame.paragraphs:
            paragraph.font.size = Pt(48)

    # Saving file
    root.save(outputFileName)
    print("PPT Created !")

def PPTtoIMG(inputFileName, outputPath):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    deck = powerpoint.Presentations.Open(inputFileName)
    powerpoint.ActivePresentation.Export(outputPath, "JPG")

    deck.Close()
    powerpoint.Quit()
    print("Images Created")

working_dir = os.getcwd()
content = [["1. Introduction", "List 1", "List 2", "List 3"], ["2. Discussion", "List 1", "List 2", "List 3"], ["3. Conclusion", "List 1", "List 2", "List 3"]]
createPresentation(working_dir + "\MyPresentation.pptx", content, "Topic")
# PPTtoIMG(working_dir +"\MyPresentation.pptx", working_dir +"\Images")