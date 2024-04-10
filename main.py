import os

import comtypes.client
import pythoncom
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from wtforms import Form, StringField, validators, SubmitField
import aicontent
#from Demo import *
from pptx.enum.text import PP_ALIGN


app = Flask(__name__)

class QueryForm(Form):
    mainpoints = StringField("",validators=[validators.input_required()])
    topic = StringField("", validators=[validators.input_required()])
    bg = StringField("")
    font = StringField("")
    submit = SubmitField("Submit")
    template_1 = SubmitField()
    template_2 = SubmitField()
    template_3 = SubmitField()

def deleteFile(dir):
    for root, dirs, files in os.walk(dir):
        for file in files:
            os.remove(os.path.join(root, file))


def createPresentation(outputFileName, content, topic, mainpoints, font_color, bg_color):
    deleteFile("C:/Users/faruk/PycharmProjects/presentations/output")
    root = Presentation()
    for slide in root.slides:
        root.slides.delete(slide)


    first_slide_layout = root.slide_layouts[0]
    second_slide_layout = root.slide_layouts[1]
    """ Ref for slide types:
    0 -> title and subtitle
    1 -> title and content
    2 -> section header
    3 -> two content
    4 -> Comparison
    5 -> Title only
    6 -> Blank
    7 -> Content with caption
    8 -> Pic with caption
    """
    # Creating slide object to add
    # in ppt i.e. Attaching slides
    # with Presentation i.e. ppt
    title_slide = root.slides.add_slide(first_slide_layout)
    background = title_slide.background
    fill = background.fill
    fill.solid()
    if bg_color is None:
       fill.fore_color.rgb = RGBColor.from_string("100E1B")
    else:
        fill.fore_color.rgb = RGBColor.from_string(bg_color)
    title_slide.placeholders[0].text= f'{topic}'
    for paragraph in title_slide.placeholders[0].text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.name = "Arial Black"
        if font_color is None:
           paragraph.font.color.rgb = RGBColor.from_string("efefef")
        else:
           paragraph.font.color.rgb = RGBColor.from_string(font_color)
        paragraph.margin_left = Inches(5)
    #content_slide
    content_slide = root.slides.add_slide(second_slide_layout)
    background = content_slide.background
    fill = background.fill
    fill.solid()
    if bg_color is None:
       fill.fore_color.rgb = RGBColor.from_string("100E1B")
    else:
        fill.fore_color.rgb = RGBColor.from_string(bg_color)
    content_slide.shapes.title.text = "Content"
    for paragraph in content_slide.placeholders[0].text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.name = "Arial Black"
        if font_color is None:
           paragraph.font.color.rgb = RGBColor.from_string("efefef")
        else:
           paragraph.font.color.rgb = RGBColor.from_string(font_color)

    content_slide.placeholders[1].text = mainpoints
    for paragraph in content_slide.placeholders[1].text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.name = "Arial Black"
        if font_color is None:
           paragraph.font.color.rgb = RGBColor.from_string("efefef")
        else:
           paragraph.font.color.rgb = RGBColor.from_string(font_color)

    for row in content:
        slide = root.slides.add_slide(second_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        if bg_color is None:
           fill.fore_color.rgb = RGBColor.from_string("100E1B")
        else:
           fill.fore_color.rgb = RGBColor.from_string(bg_color)
        # Adding title and subtitle in
        # slide i.e. first page of slide
        slide.placeholders[0].text = row[0]
        for paragraph in slide.placeholders[0].text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.name = "Arial Black"
            if font_color is None:
               paragraph.font.color.rgb = RGBColor.from_string("efefef")
            else:
               paragraph.font.color.rgb = RGBColor.from_string(font_color)
        # We have different formats of
        # subtitles in ppts, for simple
        # subtitle this method should
        # implemented, you can change
        # 0 to 1 for different design
        text = ""
        check = True
        for col in row:
            if check:
                check = False
            else:
                text = text + col + "\n"

        slide.placeholders[1].text = text
        for paragraph in slide.placeholders[1].text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.name = "Arial Black"
            if font_color is None:
               paragraph.font.color.rgb = RGBColor.from_string("efefef")
            else:
               paragraph.font.color.rgb = RGBColor.from_string(font_color)
    #conclusion slide
    last_slide = root.slides.add_slide(first_slide_layout)
    background = last_slide.background
    fill = background.fill
    fill.solid()
    if bg_color is None:
       fill.fore_color.rgb = RGBColor.from_string("100E1B")
    else:
        fill.fore_color.rgb = RGBColor.from_string(bg_color)
    last_slide.placeholders[0].text= 'Thanks for Listening'
    for paragraph in last_slide.placeholders[0].text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.name = "Arial Black"
        if bg_color is None:
           paragraph.font.color.rgb = RGBColor.from_string("efefef")
        else:
           paragraph.font.color.rgb = RGBColor.from_string(font_color)
        paragraph.margin_left = Inches(5)
    # Saving file
    root.save(outputFileName)
    print("ppt created")


def PPTtoIMG(inputFileName, outputPath):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application",pythoncom.CoInitialize())
    powerpoint.Visible = 1

    deck = powerpoint.Presentations.Open(inputFileName)
    powerpoint.ActivePresentation.Export(outputPath, "JPG")

    deck.Close()
    powerpoint.Quit()
    print("Images Created")

def listFiles(dir):
    file_list = []
    fileNames = os.listdir(dir)
    for fileName in fileNames:
        file_list.append(f"Images/{fileName}")
    print(file_list)
    return file_list





@app.route("/", methods=["GET", "POST"])
def home():
    FOLDER_PATH = "C:/Users/faruk/PycharmProjects/presentations/static/Images"
    form = QueryForm()
    deleteFile(FOLDER_PATH)
    if request.method == "POST":
        try:
            p_mainpoints = request.form.get("mainpoints")
            p_mainpoints = p_mainpoints.replace(",","\n")
            p_font = request.form.get("font")
            p_bg = request.form.get("bg")
            print(p_mainpoints)
            p_topic = request.form.get("topic")
            templates_pp = request.form.get("templates")
            print(templates_pp)
            chosen_temp = ""
            if templates_pp == "template_1":
               chosen_temp = "Management Consulting Toolkit by Slidesgo.pptx"
            elif templates_pp == "template_2":
               chosen_temp = "Basic Template_ Acid Duotone Theme for Marketing by Slidesgo.pptx"
            p_content = aicontent.produceQuery(p_mainpoints, p_topic)
            splitted_content = p_content.split("#")
            slide_content_unedited = []
            for page in splitted_content:
                bullet_points = page.split("~")
                slide_content_unedited.append(bullet_points)
            print(slide_content_unedited)
            slide_content = []
            slide_content_unedited.pop(0)
            for page in slide_content_unedited:
                slide = []
                page[1::] = [''.join(page[2::])]
                slide.append("".join(page))
                slide_content.append(slide)
            print(slide_content)

            working_dir = os.getcwd()
            createPresentation("output\my_presentation.pptx",slide_content_unedited,p_topic, p_mainpoints, p_font,p_bg)
            #create_pres(p_topic, p_mainpoints, slide_content, "template 3.pptx")
            PPTtoIMG(working_dir + "\output\my_presentation.pptx", working_dir + "\static\Images")
            slides = listFiles(FOLDER_PATH)
            title_slide = slides[0]
            slides.pop(0)
            return render_template("index.html", form=form, content=splitted_content, points=p_mainpoints, slides=slides,title_slide= title_slide, display=True, topic=p_topic,error = False)
        except:
            return render_template("index.html", error = True)

    return render_template("index.html", form=form, display=False)

@app.route("/download")
def download_file():
    p = "Mypresentation.pptx"
    return send_file(p, as_attachment=True)



if __name__ == '__main__':
    app.run(debug=True)
